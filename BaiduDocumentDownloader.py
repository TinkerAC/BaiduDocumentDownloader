import json
import os
import random
import re
import sys
import time
from tqdm import tqdm
from datetime import datetime

import requests
from bs4 import BeautifulSoup
from docx import Document
from docx.oxml.ns import qn
from docx.shared import Pt
from fake_useragent import UserAgent
from pptx import Presentation
from selenium import webdriver
from selenium.webdriver.chrome.options import Options
from selenium.webdriver.common.by import By
from selenium.webdriver.support import expected_conditions as EC
from selenium.webdriver.support.ui import WebDriverWait


class BaiduDocumentDownloader:

    def __init__(self):
        ua = UserAgent()
        self.headers = {'User-Agent': ua.random,
                        "Referer": "https://wenku.baidu.com"
                        }
        self.WORK_PATH = os.path.dirname(os.path.abspath(__file__))
        self.PATTERN_DOC_ID = r"/[\d\w]{40}|/[\d\w]{24}"

    @staticmethod
    def _validate_doc_id(doc_id):

        return re.match(self.PATTERN_DOC_ID, doc_id)

    @staticmethod
    def _save_json(json_struct, name):
        with open(name, "w", encoding='utf-8') as file:
            json.dump(json_struct, file, ensure_ascii=False)
        # print(f"{name}.json已保存！")

    @staticmethod
    def _prase_file_type(soup):
        pattern = r'"fileType":"(.*?)"'
        script_tags = soup.findAll("script")
        for tag in script_tags:
            matched = re.search(pattern, tag.text)
            if matched and matched.group(1):
                return matched.group(1)
        print("自动获取文档类型失败!")
        type_ = input(r"请手动输入文档类型:(word\pdf\excel\ppt\txt)")
        return type_

    @staticmethod
    def _download_txt(doc_id, path):
        url = f"https://wenku.baidu.com/view/{doc_id}"
        response = requests.get(url)
        soup = BeautifulSoup(response.text, "html.parser")
        p_list = soup.findAll("p", class_="p-txt")
        with open(path, "w") as file:
            for p in p_list:
                file.write(p.text)

        print("文档已成功保存到指定目录：", path)

    @staticmethod
    def _prase_title(soup):
        return soup.find('title').text.replace('- 百度文库', "").strip()

    @staticmethod
    def _create_pptx(image_folder, pptx_file):
        prs = Presentation()

        slide_width = prs.slide_width  # 获取幻灯片宽度
        slide_height = prs.slide_height  # 获取幻灯片高度

        for image_file in sorted(os.listdir(image_folder)):
            if image_file.endswith((".png", ".jpg", ".jpeg")):
                slide = prs.slides.add_slide(prs.slide_layouts[5])  # 使用空白幻灯片布局
                pic = slide.shapes.add_picture(os.path.join(image_folder, image_file), 0, 0, width=slide_width,
                                               height=slide_height)

        prs.save(pptx_file)

    @staticmethod
    def _download_images(url_list, download_folder_path):

        count = 0
        for url in tqdm(url_list):
            file_name = f"{count}.jpg"
            save_path = os.path.join(download_folder_path, file_name)

            response = requests.get(url, stream=True)
            if response.status_code == 200 or 206:
                with open(save_path, 'wb') as file:
                    for chunk in response.iter_content(chunk_size=8192):
                        file.write(chunk)

            else:
                print("下载失败，HTTP状态码:", response.status_code)


    @staticmethod
    def _get_soup_with_driver(url_main):
        # 初始化Chrome浏览器的无头模式
        chrome_options = Options()
        chrome_options.add_argument("--headless")  # 设置为无头模式
        chrome_options.add_argument("--disable-gpu")

        driver = webdriver.Chrome(options=chrome_options)

        try:
            # 访问文档主页
            driver.get(url_main)

            # 等待页面标题加载完成
            WebDriverWait(driver, 10).until(
                EC.presence_of_element_located((By.TAG_NAME, "title"))
            )
            page_source = driver.page_source
        except TimeoutError:
            print("页面加载超时！将无返回值返回。")
            page_source = None
        finally:
            driver.quit()

        return page_source

    def download(self, url_or_id=None):
        if not url_or_id:
            while True:
                doc_id = input('请输入文档id/URL:')
                if doc_id:
                    break

        else:
            doc_id = url_or_id

        # 从URL中提取文档ID
        matched = re.search(self.PATTERN_DOC_ID, doc_id)
        if matched:
            doc_id = matched.group()[1:]
            input(f"请检查文档ID:{doc_id}")

        current_timestamp = int(time.time())

        url_main = f'https://wenku.baidu.com/view/{doc_id}.html?_wkts_={current_timestamp}&fr=hp_Database'

        response_main = requests.get(url_main)

        if response_main.status_code != 200:
            print(f'获取主页面信息失败,HTTP状态码{response_main.status_code}')
            print("正在使用Chrome WebDriver...")
            page_source = self._get_soup_with_driver(url_main)
        else:
            page_source = response_main.text

        soup_main = BeautifulSoup(page_source, "html.parser")

        title = self._prase_title(soup_main)

        all_downloaded_file_path = os.path.join(self.WORK_PATH, "已下载文档")
        os.makedirs(all_downloaded_file_path, exist_ok=True)  # 创建包含所有已下载文件的文件夹

        download_folder_path = os.path.join(all_downloaded_file_path, title)
        os.makedirs(download_folder_path, exist_ok=True)  # 创建当前下载的文件的文件夹

        url_data = f"https://wenku.baidu.com/ndocview/readerinfo?docId={doc_id}&pn=100"  # 调用数据API
        response = requests.get(url_data, headers=self.headers)
        print(url_data)
        data_url = json.loads(response.text)

        self._save_json(data_url, "data_url.json")

        if response.status_code == 200:
            print('响应状态:', response.status_code, '正常')
        else:
            quit()

        # 文档可用性检查
        api_status_code = data_url["status"]['code']
        if api_status_code == 0:
            print(f'接口响应代码:{api_status_code}   正常')
        else:
            print(f"data_url:{data_url}")
            raise Exception(f'接口响应代码:{api_status_code}   异常')

        all_page = data_url["data"]["freePage"]
        free_page = data_url['data']["showPage"]
        print(f"文档可用性：{free_page}页/{all_page}页")

        file_type = self._prase_file_type(soup_main)
        print(f"文档类型:{file_type}")

        if file_type in ["pdf", 'word', "excel"]:
            text = ""
            content_row = []  # 用于存储每行文字的列表
            json_data = None
            for i in data_url["data"]['htmlUrls']['json']:
                url = i["pageLoadUrl"]
                response = requests.get(url, headers=self.headers)
                content = response.text
                pattern = r'\{.*\}'  # 匹配以 '{' 开头，以 '}' 结尾的部分
                matched = re.search(pattern, content)
                json_data = json.loads(matched.group())

                last_y = ""
                count = 0
                for j in json_data["body"]:  # json_data["body"]是每一页的内容的body

                    y = j["p"]["y"]

                    if y != last_y and text:
                        last_y = y
                        count += 1
                        content_row.append(text)
                        text = ""

                    try:
                        text += j['c']
                    except:  # 有可能为一个dict,待处理

                        pass

            self._save_json(json_data, "url1.json")

            text = ""
            for i in content_row:
                text += i + "\n"

            print(f'《{title}》，共计{free_page}页。')

            # 创建一个新的Word文档
            doc = Document()

            # 添加文本到文档
            paragraph = doc.add_paragraph(text)

            # 设置字体样式为宋体
            run = paragraph.runs[0]
            run.font.name = '宋体'

            # 为确保在Word中也是宋体，需要设置字体的ascii和eastAsia属性
            run.font.name = 'Times New Roman'
            run._element.rPr.rFonts.set(qn('w:eastAsia'), '宋体')

            # 可以选择设置字体大小
            run.font.size = Pt(12)

            if not title.isspace():
                file_name = f'{title}.docx'
            else:
                current_time = datetime.now().strftime("%Y%m%d%H%M%S")
                file_name = f'{current_time}.docx'

            save_path = os.path.join(download_folder_path, file_name)
            doc.save(save_path)
            img_object_list = data_url["data"]['htmlUrls']["png"]

            if img_object_list:  # 如果文档有配图,一并下载
                download_folder_path_img = os.path.join(download_folder_path, "相关图片")
                os.makedirs(download_folder_path_img, exist_ok=True)
                url_list = [obj["pageLoadUrl"] for obj in img_object_list]
                self._download_images(url_list, download_folder_path_img)
            print("\n文档已成功保存到指定目录：", download_folder_path)

        elif file_type == "ppt":
            img_urls = data_url["data"]['htmlUrls']
            self._download_images(img_urls, download_folder_path)

            image_folder_path = download_folder_path  # 替换为您的图片文件夹路径
            pptx_file_path = f'{title}.pptx'  # 指定要创建的 PowerPoint 文件名
            self._create_pptx(image_folder_path, pptx_file_path)
            print("\n文档已成功保存到指定目录：", download_folder_path)

        elif file_type == 'txt':
            file_name = f"{title}.txt"
            save_path = os.path.join(download_folder_path, file_name)
            self._download_txt(doc_id, save_path)
            print("文档已成功保存到指定目录：", download_folder_path)


if __name__ == '__main__':
    downloader = BaiduDocumentDownloader()

    docx_url = r"https://wenku.baidu.com/view/ed7db54af211f18583d049649b6648d7c0c70813.html?fr=hp_doclist&_wkts_=1713013169010"
    pptx_url = r"https://wenku.baidu.com/view/50fceb89eb7101f69e3143323968011ca200f719.html?fr=hp_doclist&_wkts_=1713013060941"
    pdf_url = r"https://wenku.baidu.com/view/78df162df505cc1755270722192e453610665bb1.html?fr=hp_doclist&_wkts_=1713013121798"
    excel_url = "https://wenku.baidu.com/view/988dc659a300a6c30d229f2e.html?fr=search-rec-1&_wkts_=1713013302604&wkQuery=excel"
    txt_url = r"https://wenku.baidu.com/view/fd8e8473f242336c1eb95ea2.html?fr=income1-doc-search&_wkts_=1713013230980&wkQuery=txt"
    # test_unit = [docx_url, pptx_url, pdf_url, excel_url, txt_url]
    # random.shuffle(test_unit)
    # for url in test_unit:
    #     print(f"正在下载:{url}")
    #     downloader.download(url)
    downloader.download(pptx_url)